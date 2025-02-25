from flask import Flask, render_template, request, jsonify
from flask_socketio import SocketIO, emit
from models import db, Order, OrderItem
from datetime import datetime, timedelta
import uuid
from dotenv import load_dotenv
import os
import json
from typing import List
from langchain.document_loaders import TextLoader, PyPDFLoader
from langchain.text_splitter import CharacterTextSplitter
from langchain.embeddings import OpenAIEmbeddings
from langchain.vectorstores import FAISS
from langchain.document_loaders.base import BaseLoader
from langchain.docstore.document import Document
from pptx import Presentation
import openai

# 加載環境變數
load_dotenv()

app = Flask(__name__)
socketio = SocketIO(app, cors_allowed_origins="*")

# 配置 SQLAlchemy
app.config['SQLALCHEMY_DATABASE_URI'] = 'sqlite:///orders.db'
app.config['SQLALCHEMY_TRACK_MODIFICATIONS'] = False
db.init_app(app)

# 初始化 OpenAI
openai.api_key = os.getenv('OPENAI_API_KEY')

# 存儲系統提示詞
system_prompt = """
你是一位熱情友善的餐廳服務員，名叫小晴。你的回答風格應該永遠保持親切和專業。

當顧客詢問菜單時，請使用以下格式回答：

### 主食系列

| 品名 | 價格 | 客製化 | 備註 |
|:-----|:-----|:-------|:-----|
| 和牛肉眼牛排 | $85 | ✓ | 配以蒜香牛油和松露醬汁 |
| 龍蝦熱奶油 | $60 | ✓ | 奶油龍蝦配奶油白酒醬，灑上綠色馬蘭茲 |
| 嫩煎德國豬 | $38 | ✓ | 香煎豬腰脊配蘋果醬 |
| 鮮蝦意粉 | $28 | ✓ | 鮮蝦與意大利麵絲炒製 |
| 芋頭奧索布料 | $48 | ✓ | 紅酒醬汁淋芋頭奧索布料，搭配黑松露風味醬汁 |
| 黑松露燉飯 | $40 | ✓ | 有著黑松露香和帕瑪森起司的奶香燉飯 |

### 飲品系列

| 品名 | 價格 | 客製化 | 備註 |
|:-----|:-----|:-------|:-----|
| 招牌奶茶 | $50 | ✓ | 完美調配的奶茶 |
| 黑糖珍珠鮮奶 | $65 | ✓ | 香醇黑糖搭配鮮奶 |
| 芋頭鮮奶 | $60 | ✓ | 使用真實芋頭製作 |

溫馨提醒：

主食客製化選項：
- 主食熟度：三分、五分、七分、全熟
- 醬汁選擇：黑胡椒、蘑菇、紅酒
- 配菜選擇：時蔬、薯條、沙拉

飲品客製化選項：
- 冰量：全冰、少冰、微冰、去冰
- 糖度：全糖、少糖、半糖、微糖、無糖
- 配料：珍珠、椰果、布丁（每項 +$10）
- 杯型：中杯、大杯（大杯 +$10）

當顧客要點餐時，請確認以下信息：
1. 主食：
   - 熟度偏好
   - 醬汁選擇
   - 配菜選擇
2. 飲品：
   - 冰量
   - 糖度
   - 配料
   - 杯型

請以自然、正面的語氣與顧客互動，適時推薦特色餐點和飲品組合。

在顧客完成點餐後，請使用以下格式整理並確認訂單：

=== 您的訂單明細 ===
【主食】
- 品項：[品名]
  數量：[數量]
  價格：$[單價]
  客製：[熟度]、[醬汁]、[配菜]

【飲品】
- 品項：[品名]
  數量：[數量]
  價格：$[單價]
  客製：[冰量]、[糖度]、[配料]

=== 總計金額：$[總價] ===

請問訂單內容是否正確？如需修改請告訴我。
"""

# 存儲向量數據庫
vector_store = None

@app.route('/')
def home():
    return render_template('index.html')

@app.route('/admin')
def admin():
    return render_template('admin.html')

@app.route('/reports')
def reports():
    return render_template('reports.html')

@app.route('/api/reports')
def get_reports():
    report_type = request.args.get('type', 'daily')
    start_date = request.args.get('start')
    end_date = request.args.get('end')
    
    # 將字符串轉換為日期對象
    start_date = datetime.strptime(start_date, '%Y-%m-%d').date()
    end_date = datetime.strptime(end_date, '%Y-%m-%d').date()
    
    # 查詢指定日期範圍內的訂單
    orders = Order.query.filter(
        db.func.date(Order.created_at) >= start_date,
        db.func.date(Order.created_at) <= end_date
    ).all()
    
    # 計算總營業額和訂單數
    total_revenue = sum(order.total_amount for order in orders)
    total_orders = len(orders)
    average_order = total_revenue / total_orders if total_orders > 0 else 0
    
    # 計算銷售趨勢
    trends = []
    current_date = start_date
    
    if report_type == 'monthly':
        # 按月统计
        while current_date <= end_date:
            # 获取当月最后一天
            if current_date.month == 12:
                month_end = datetime(current_date.year + 1, 1, 1).date() - timedelta(days=1)
            else:
                month_end = datetime(current_date.year, current_date.month + 1, 1).date() - timedelta(days=1)
            
            month_end = min(month_end, end_date)
            month_orders = [o for o in orders if current_date <= o.created_at.date() <= month_end]
            month_revenue = sum(o.total_amount for o in month_orders)
            
            trends.append({
                'date': f"{current_date.strftime('%Y-%m')} ({current_date.strftime('%Y-%m-%d')} ~ {month_end.strftime('%Y-%m-%d')})",
                'revenue': month_revenue,
                'orders': len(month_orders)
            })
            
            # 移到下个月第一天
            if current_date.month == 12:
                current_date = datetime(current_date.year + 1, 1, 1).date()
            else:
                current_date = datetime(current_date.year, current_date.month + 1, 1).date()
                
    elif report_type == 'weekly':
        # 調整到週一開始
        while current_date.weekday() != 0:  # 0 = Monday
            current_date -= timedelta(days=1)
        
        # 按週統計
        while current_date <= end_date:
            week_end = min(current_date + timedelta(days=6), end_date)
            week_orders = [o for o in orders if current_date <= o.created_at.date() <= week_end]
            week_revenue = sum(o.total_amount for o in week_orders)
            trends.append({
                'date': f"{current_date.strftime('%Y-%m-%d')} ~ {week_end.strftime('%Y-%m-%d')}",
                'revenue': week_revenue,
                'orders': len(week_orders)
            })
            current_date += timedelta(days=7)
    else:
        # 按日統計
        while current_date <= end_date:
            day_orders = [o for o in orders if o.created_at.date() == current_date]
            day_revenue = sum(o.total_amount for o in day_orders)
            trends.append({
                'date': current_date.strftime('%Y-%m-%d'),
                'revenue': day_revenue,
                'orders': len(day_orders)
            })
            current_date += timedelta(days=1)
    
    # 計算熱門商品
    item_stats = {}
    for order in orders:
        for item in order.items:
            if item.item_name not in item_stats:
                item_stats[item.item_name] = {'quantity': 0, 'revenue': 0}
            item_stats[item.item_name]['quantity'] += item.quantity
            item_stats[item.item_name]['revenue'] += item.price * item.quantity
    
    popular_items = [
        {'name': name, 'quantity': stats['quantity'], 'revenue': stats['revenue']}
        for name, stats in item_stats.items()
    ]
    popular_items.sort(key=lambda x: x['quantity'], reverse=True)
    
    # 準備詳細數據
    details = [{
        'date': trend['date'],
        'orders': trend['orders'],
        'revenue': trend['revenue'],
        'average': trend['revenue'] / trend['orders'] if trend['orders'] > 0 else 0
    } for trend in trends]
    
    return jsonify({
        'summary': {
            'total_revenue': total_revenue,
            'total_orders': total_orders,
            'average_order': average_order
        },
        'trends': trends,
        'popular_items': popular_items[:10],  # 只返回前10個熱門商品
        'details': details
    })

@app.route('/orders')
def orders():
    return render_template('orders.html')

@app.route('/api/orders')
def get_orders():
    # 獲取篩選參數
    status = request.args.get('status')
    date = request.args.get('date')
    search = request.args.get('search')
    
    # 構建查詢
    query = Order.query
    
    if status:
        query = query.filter(Order.status == status)
    if date:
        query = query.filter(db.func.date(Order.created_at) == date)
    if search:
        query = query.filter(
            db.or_(
                Order.order_number.contains(search),
                Order.customer_name.contains(search)
            )
        )
    
    orders = query.order_by(Order.created_at.desc()).all()
    return jsonify({
        'orders': [{
            'id': order.id,
            'order_number': order.order_number,
            'customer_name': order.customer_name,
            'total_amount': order.total_amount,
            'status': order.status,
            'created_at': order.created_at.isoformat()
        } for order in orders]
    })

@app.route('/api/orders/<int:order_id>')
def get_order(order_id):
    order = Order.query.get_or_404(order_id)
    return jsonify({
        'id': order.id,
        'order_number': order.order_number,
        'customer_name': order.customer_name,
        'total_amount': order.total_amount,
        'status': order.status,
        'created_at': order.created_at.isoformat(),
        'items': [{
            'item_name': item.item_name,
            'quantity': item.quantity,
            'price': item.price,
            'options': item.options
        } for item in order.items]
    })

@app.route('/api/orders/<int:order_id>', methods=['PATCH'])
def update_order(order_id):
    order = Order.query.get_or_404(order_id)
    data = request.json
    
    if 'status' in data:
        order.status = data['status']
        db.session.commit()
    
    return jsonify({'message': 'Order updated successfully'})

@app.route('/api/orders', methods=['POST'])
def create_order():
    data = request.json
    
    # 創建訂單
    order = Order(
        order_number=f"ORD{datetime.now().strftime('%Y%m%d')}{str(uuid.uuid4())[:8]}",
        customer_name=data.get('customer_name', 'Anonymous'),
        total_amount=sum(item['price'] * item['quantity'] for item in data['items'])
    )
    
    # 添加訂單項目
    for item_data in data['items']:
        item = OrderItem(
            item_name=item_data['item_name'],
            quantity=item_data['quantity'],
            price=item_data['price'],
            options=item_data.get('options', {})
        )
        order.items.append(item)
    
    db.session.add(order)
    db.session.commit()
    
    return jsonify({
        'message': 'Order created successfully',
        'order_id': order.id,
        'order_number': order.order_number
    })

@app.route('/chat', methods=['POST'])
def chat():
    try:
        data = request.get_json()
        if not data or 'message' not in data:
            return jsonify({"error": "缺少訊息內容"}), 400

        user_message = data['message'].strip()
        if not user_message:
            return jsonify({"error": "訊息不能為空"}), 400

        # 如果有向量存儲，先進行相關信息檢索
        context = ""
        if vector_store:
            try:
                docs = vector_store.similarity_search(user_message, k=2)
                
                # 格式化文檔內容為 Markdown 表格
                formatted_content = []
                current_category = ''
                table_rows = []
                
                for doc in docs:
                    lines = doc.page_content.split('\n')
                    for line in lines:
                        line = line.strip()
                        if not line:  # 跳過空行
                            continue
                            
                        if line.startswith('###'):
                            # 如果有之前的表格，先添加到格式化內容中
                            if table_rows:
                                formatted_content.extend([
                                    '',
                                    '| 品名 | 價格 | 客製化 | 備註 |',
                                    '|------|------|--------|------|',
                                    *table_rows,
                                    ''
                                ])
                                table_rows = []
                            
                            current_category = line.strip('# ')
                            formatted_content.extend([
                                '',
                                f'### {current_category}',
                                ''
                            ])
                        elif line.startswith('-'):
                            # 處理菜單項目
                            item = line.strip('- ').strip()
                            if ' - ' in item:  # 使用 - 作為分隔符
                                parts = item.split(' - ')
                                name = parts[0].strip()
                                description = ' - '.join(parts[1:]).strip()
                                
                                # 從描述中提取價格和備註
                                if '**$' in description:
                                    price_parts = description.split('**$')
                                    note = price_parts[0].strip().strip('**')
                                    price = f'${price_parts[1].strip().strip("**")}'
                                else:
                                    note = description
                                    price = 'N/A'
                                
                                # 根據產品類型添加客製化選項
                                customization = '✔️' if any(keyword in name.lower() for keyword in ['奶茶', '茶', '哈哈', '咖啡']) else '-'
                                
                                table_rows.append(f'| {name} | {price} | {customization} | {note} |')
                            else:
                                formatted_content.append(f'- {item}')
                        else:
                            formatted_content.append(line)
                
                # 添加最後一個表格
                if table_rows:
                    formatted_content.extend([
                        '',
                        '| 品名 | 價格 | 客製化 | 備註 |',
                        '|------|------|--------|------|',
                        *table_rows,
                        '',
                        '\n備註：',
                        '- ✔️ 表示可以調整冰量和糖度',
                        '- 可加珍珠、椰果、布丁等配料（每項 +$10）',
                        '- 大杯升級 +$10',
                        ''
                    ])
                
                # 添加最后一個表格
                if table_rows:
                    formatted_content.append('| 品名 | 價格 | 備註 |')
                    formatted_content.append('|------|------|------|')
                    formatted_content.extend(table_rows)
                
                context = "\n".join(formatted_content)
            except Exception as e:
                print(f"知識庫查詢錯誤：{str(e)}")

        # 準備完整的系統提示詞
        base_prompt = system_prompt or "你是一位有趣的飲料店服務員，名叫小晴。你會以有趣的方式詢問顧客需求，並推薦招牌飲品。如果顧客要下單，請以 JSON 格式輸出訂單資訊。"
        
        # 添加訂單相關的提示詞
        order_prompt = """
如果顧客要下單，請以下列 JSON 格式輸出訂單資訊：
```json
{
    "is_order": true,
    "items": [
        {
            "item_name": "商品名稱",
            "quantity": 數量,
            "price": 單價,
            "options": {
                "ice": "冰量",
                "sugar": "糖量"
            }
        }
    ],
    "total_amount": 總價
}
```
"""

        full_system_prompt = (
            f"{base_prompt}\n{order_prompt}\n\n相關商品資訊：\n{context}" 
            if context else f"{base_prompt}\n{order_prompt}"
        )

        # 調用 OpenAI API
        try:
            response = openai.ChatCompletion.create(
                model="gpt-3.5-turbo",
                messages=[
                    {"role": "system", "content": full_system_prompt},
                    {"role": "user", "content": user_message}
                ],
                temperature=0.7,  # 增加一些變化性，讓回答更有趣
                max_tokens=800   # 增加回答長度以容納訂單資訊
            )
            
            response_text = response.choices[0].message.content
            
            # 嘗試尋找並解析 JSON 訂單
            try:
                # 尋找 JSON 內容
                import re
                # 先尋找在 ```json ... ``` 內的 JSON
                json_match = re.search(r'```json\s*(\{[^`]*\})\s*```', response_text)
                if json_match:
                    json_text = json_match.group(1)
                else:
                    # 如果沒有找到，則直接尋找 JSON 格式
                    json_match = re.search(r'\{[^{]*"is_order"\s*:\s*true[^}]*\}', response_text)
                    if json_match:
                        json_text = json_match.group(0)
                    else:
                        json_text = None
                
                if json_text:
                    order_data = json.loads(json_text)
                    # 如果是訂單，檢查是否有進行中的訂單
                    if order_data.get('is_order'):
                        # 查找最近的未完成訂單
                        recent_order = Order.query.filter_by(status='pending').order_by(Order.created_at.desc()).first()
                        
                        if recent_order and (datetime.now() - recent_order.created_at).total_seconds() < 600:  # 10分鐘內的訂單
                            # 追加到現有訂單
                            order = recent_order
                            # 更新總金額
                            order.total_amount += order_data['total_amount']
                        else:
                            # 創建新訂單
                            order = Order(
                                order_number=str(uuid.uuid4())[:8],
                                total_amount=order_data['total_amount'],
                                status='pending',
                                created_at=datetime.now()
                            )
                            db.session.add(order)
                            db.session.flush()  # 確保取得 order.id
                        
                        # 添加新的訂單項目
                        for item in order_data['items']:
                            order_item = OrderItem(
                                order_id=order.id,
                                item_name=item['item_name'],
                                quantity=item['quantity'],
                                price=item['price'],
                                options=json.dumps(item.get('options', {}))
                            )
                            db.session.add(order_item)
                        
                        db.session.commit()
                        
                        # 獲取完整的訂單明細
                        all_items = OrderItem.query.filter_by(order_id=order.id).all()
                        items_summary = "\n=== 當前訂單明細 ===\n"
                        total = 0
                        for item in all_items:
                            options = json.loads(item.options)
                            items_summary += f"- {item.item_name}\n"
                            items_summary += f"  * 數量：{item.quantity}\n"
                            items_summary += f"  * 單價：{item.price} 元\n"
                            items_summary += f"  * 客製化：冰量 {options.get('ice', '正常')}、糖量 {options.get('sugar', '正常')}\n"
                            total += item.quantity * item.price
                        items_summary += f"=== 總計：{total} 元 ===\n"
                        items_summary += f"訂單編號：{order.order_number}\n"
                        
                        # 將訂單明細加入回應中
                        try:
                            # 如果response_text是完整的JSON，直接更新
                            try:
                                response_dict = json.loads(response_text)
                                response_dict.update({
                                    "order_summary": items_summary,
                                    "order_number": order.order_number
                                })
                                response_text = json.dumps(response_dict, ensure_ascii=False)
                            except json.JSONDecodeError:
                                # 如果response_text不是完整的JSON，构造一个新的响应
                                response_text = json.dumps({
                                    "response": response_text,
                                    "order_summary": items_summary,
                                    "order_number": order.order_number
                                }, ensure_ascii=False)
                        except json.JSONDecodeError:
                            escaped_summary = items_summary.replace('\n', '\\n').replace('"', '\\"')
                            response_text = response_text.rstrip('}') + \
                                f", \"order_summary\": \"{escaped_summary}\", \"order_number\": \"{order.order_number}\"}}"
            except Exception as e:
                print(f"處理訂單時發生錯誤：{str(e)}")
                # 繼續處理，不中斷對話
            
            # 如果response_text已经是JSON字符串，解析它
            try:
                response_dict = json.loads(response_text)
                response_dict["context_used"] = bool(context)
                return jsonify(response_dict)
            except json.JSONDecodeError:
                # 如果不是JSON，使用原始格式
                return jsonify({
                    "response": response_text,
                    "context_used": bool(context)
                })
            
        except openai.error.OpenAIError as e:
            print(f"OpenAI API 錯誤：{str(e)}")
            return jsonify({"error": "處理您的訊息時發生錯誤，請稍後再試"}), 500
            
    except Exception as e:
        print(f"聊天功能錯誤：{str(e)}")
        return jsonify({"error": "系統發生錯誤，請稍後再試"}), 500

class PowerPointLoader(BaseLoader):
    """加載 PowerPoint 文件的自定義加載器"""
    
    def __init__(self, file_path: str):
        self.file_path = file_path
    
    def load(self) -> List[Document]:
        prs = Presentation(self.file_path)
        documents = []
        
        for i, slide in enumerate(prs.slides, 1):
            # 提取幻燈片中的所有文本
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = shape.text.strip()
                    if text:
                        texts.append(text)
            
            if texts:
                # 將幻燈片的所有文本合併
                text_content = "\n".join(texts)
                # 創建文檔對象，包含幻燈片編號的元數據
                doc = Document(
                    page_content=text_content,
                    metadata={"source": self.file_path, "slide_number": i}
                )
                documents.append(doc)
        
        return documents

@app.route('/admin/upload', methods=['POST'])
def upload_document():
    print("接收到文件上传请求")
    print(f"请求头: {dict(request.headers)}")
    
    if 'file' not in request.files:
        print("错误：没有文件字段")
        return jsonify({"error": "沒有提供文件"}), 400

    file = request.files['file']
    print(f"文件名: {file.filename}")
    print(f"文件类型: {file.content_type}")
    
    if file.filename == '':
        print("错误：文件名为空")
        return jsonify({"error": "沒有選擇文件"}), 400

    # 检查文件大小和内容
    try:
        file_content = file.read()
        print(f"文件大小: {len(file_content)} bytes")
        print(f"文件内容前100字节: {file_content[:100]}")
        
        if len(file_content) > 10 * 1024 * 1024:  # 10MB in bytes
            print("错误：文件太大")
            return jsonify({"error": "文件大小不能超過10MB"}), 400
        
        if len(file_content) == 0:
            print("错误：文件内容为空")
            return jsonify({"error": "文件內容不能為空"}), 400

        # 重置文件指针
        file.seek(0)
        
        # 保存文件
        filename = file.filename
        file_path = os.path.join('uploads', filename)
        os.makedirs('uploads', exist_ok=True)
        print(f"准备保存文件到: {file_path}")
    except Exception as e:
        print(f"处理文件时出错: {str(e)}")
        return jsonify({"error": f"處理文件時出錯: {str(e)}"}), 500
    
    try:
        # 根據文件類型選擇加載器
        if filename.lower().endswith('.pdf'):
            loader = PyPDFLoader(file_path)
        elif filename.lower().endswith(('.pptx', '.ppt')):
            loader = PowerPointLoader(file_path)
        elif filename.lower().endswith(('.txt', '.md')):
            # 对于文本文件，检查编码
            try:
                file_content.decode('utf-8')
            except UnicodeDecodeError:
                return jsonify({"error": "文本文件必須為UTF-8編碼"}), 400
            loader = TextLoader(file_path)
        else:
            return jsonify({"error": "不支持的文件格式。請上傳 PDF、PowerPoint 或文本文件。"}), 400
            
        # 保存文件
        file.save(file_path)

        # 加載文檔
        documents = loader.load()
        
        # 分割文本
        text_splitter = CharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
        texts = text_splitter.split_documents(documents)

        # 創建向量存儲
        global vector_store
        vector_store = FAISS.from_documents(texts, OpenAIEmbeddings())

        return jsonify({
            "message": "文件上傳並處理成功",
            "details": {
                "filename": filename,
                "document_count": len(documents),
                "chunk_count": len(texts)
            }
        })
    
    except Exception as e:
        return jsonify({"error": f"處理文件時發生錯誤：{str(e)}"}), 500

@app.route('/admin/prompt', methods=['POST'])
def update_prompt():
    try:
        # 確保請求內容是 JSON
        if not request.is_json:
            return jsonify({"error": "請求必須是 JSON 格式"}), 400

        data = request.get_json()
        if not data:
            return jsonify({"error": "缺少請求內容"}), 400

        if 'prompt' not in data:
            return jsonify({"error": "缺少 prompt 參數"}), 400

        new_prompt = data['prompt']
        if new_prompt is None:
            return jsonify({"error": "prompt 不能為 null"}), 400

        new_prompt = str(new_prompt).strip()
        if not new_prompt:
            return jsonify({"error": "提示詞不能為空"}), 400

        if len(new_prompt) > 1000:
            return jsonify({"error": "提示詞長度不能超過 1000 字元"}), 400

        # 檢查不當詞彙
        inappropriate_words = [
            "凶狠", "暴躁", "殺人", "自殺", "暴力", "仓皇", "武器",
            "笨蛋", "白痴", "廢物", "懶惰", "討厭", "不耐煩", "生氣"
        ]
        found_words = [word for word in inappropriate_words if word in new_prompt]
        if found_words:
            return jsonify({
                "error": "提示詞包含不當詞彙",
                "details": f"請移除不當詞彙：{', '.join(found_words)}"
            }), 400

        # 構建完整的系統提示詞
        base_prompt = f"你是一位{new_prompt}的飲料店服務員，名叫小晴。你的主要工作是推薦和銷售飲品。"
        service_traits = """你會以親切有趣的方式詢問顧客需求，並推薦飲品。你的回答要簡短但有趣，讓顧客感到開心。

當顧客詢問飲品時，你會推薦以下特色飲品：
1. 珍珠奶茶 - NT$60，可選擇冰量和糖度
2. 黑糖珍珠鮮奶 - NT$75，店長推薦，使用紐西蘭鮮奶和手工黑糖珍珠
3. 四季春茶 - NT$45，可選擇冰量和糖度
4. 抹茶拿鈴 - NT$70，使用日本進口抹茶粉

冰量選擇：正常冰、少冰、微冰、去冰
糖度選擇：全糖、七分糖、半糖、三分糖、無糖

營業時間：
週一至週五：10:00-21:00
週六至週日：11:00-22:00"""

        order_handling = """
當顧客要下單時，請按以下格式確認訂單並輸出：

=== 訂單明細 ===
- 飲品：[名稱]
- 數量：[數量]
- 單價：[價格] 元
- 客製化：
  * 冰量：[冰量]
  * 糖量：[糖量]
=== 總計：[總金額] 元 ===

然後將訂單轉換為以下 JSON 格式：
{
    "is_order": true,
    "items": [
        {
            "item_name": "飲品名稱",
            "quantity": 數量,
            "price": 單價,
            "options": {
                "ice": "冰量",
                "sugar": "糖量"
            }
        }
    ],
    "total_amount": 總金額
}
"""

        global system_prompt
        system_prompt = f"{base_prompt} {service_traits} {order_handling}"

        return jsonify({
            "message": "提示詞更新成功",
            "current_prompt": system_prompt
        })

    except Exception as e:
        print(f"更新提示詞錯誤：{str(e)}")
        return jsonify({"error": "系統錯誤，請稍後再試"}), 500

def broadcast_report_update():
    # 获取最新报表数据
    current_date = datetime.now().date()
    start_date = current_date
    end_date = current_date
    
    orders = Order.query.filter(
        db.func.date(Order.created_at) >= start_date,
        db.func.date(Order.created_at) <= end_date
    ).all()
    
    # 计算实时数据
    total_revenue = sum(order.total_amount for order in orders)
    total_orders = len(orders)
    average_order = total_revenue / total_orders if total_orders > 0 else 0
    
    # 计算热门商品
    item_stats = {}
    for order in orders:
        for item in order.items:
            if item.item_name not in item_stats:
                item_stats[item.item_name] = {'quantity': 0, 'revenue': 0}
            item_stats[item.item_name]['quantity'] += item.quantity
            item_stats[item.item_name]['revenue'] += item.price * item.quantity
    
    popular_items = [
        {'name': name, 'quantity': stats['quantity'], 'revenue': stats['revenue']}
        for name, stats in item_stats.items()
    ]
    popular_items.sort(key=lambda x: x['quantity'], reverse=True)
    
    # 发送实时更新
    socketio.emit('report_update', {
        'summary': {
            'total_revenue': total_revenue,
            'total_orders': total_orders,
            'average_order': average_order
        },
        'popular_items': popular_items[:10]
    })

@socketio.on('connect')
def handle_connect():
    print('Client connected')
    broadcast_report_update()

@socketio.on('disconnect')
def handle_disconnect():
    print('Client disconnected')

if __name__ == '__main__':
    socketio.run(app, host='0.0.0.0', port=8080, debug=True)
